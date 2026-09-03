import io
import os
from typing import Dict, Any, List
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def format_curr_or_qty(num: float, is_value: bool, is_compact: bool = True) -> str:
    if num is None or num != num:
        return "Rp 0" if is_value else "0"
    
    if is_compact:
        abs_num = abs(num)
        if abs_num >= 1e9:
            val = num / 1e9
            return f"Rp {val:.2f} Miliar" if is_value else f"{val:.2f} M Unit"
        if abs_num >= 1e6:
            val = num / 1e6
            return f"Rp {val:.2f} Juta" if is_value else f"{val:.2f} Jt Unit"
        if abs_num >= 1e3:
            val = num / 1e3
            return f"Rp {val:.1f} Ribu" if is_value else f"{val:.1f} Rb Unit"
            
    if is_value:
        return f"Rp {num:,.0f}".replace(",", ".")
    return f"{num:,.0f}".replace(",", ".")

def format_period_label(period_str: str) -> str:
    if not period_str:
        return ""
    months = ["Januari", "Februari", "Maret", "April", "Mei", "Juni", "Juli", "Agustus", "September", "Oktober", "November", "Desember"]
    parts = period_str.split("-")
    if len(parts) == 2:
        try:
            m_idx = int(parts[1]) - 1
            if 0 <= m_idx < 12:
                return f"{months[m_idx]} {parts[0]}"
        except ValueError:
            pass
    return period_str

def style_cell(cell, text: str, font_size: int = 9, bold: bool = False, color: RGBColor = RGBColor(255, 255, 255), align: PP_ALIGN = PP_ALIGN.LEFT, bg_color: RGBColor = None):
    cell.text = str(text)
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        run.font.name = "Segoe UI"
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color

def set_slide_title(slide, title_text: str):
    """Sets the title in the template's title placeholder cleanly formatted without overlapping the logo on top right."""
    title_shape = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type in [1, 3]:  # TITLE or CENTER_TITLE
            title_shape = shape
            break
        elif shape.name and "Title" in shape.name:
            title_shape = shape
            break
            
    if title_shape and title_shape.has_text_frame:
        title_shape.left = Inches(0.4)
        title_shape.top = Inches(0.20)
        title_shape.width = Inches(6.8)
        title_shape.height = Inches(0.55)

        tf = title_shape.text_frame
        tf.word_wrap = True
        tf.text = ""
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = RGBColor(2, 132, 199)
    else:
        tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.20), Inches(6.8), Inches(0.55))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(2, 132, 199)

def create_dashboard_trend_chart(doi_trend: List[Dict[str, Any]], width: int = 1600, height: int = 1000) -> Image.Image:
    """Generates a high-DPI dark glassmorphic trend chart matching the dashboard UI styling."""
    img = Image.new("RGBA", (width, height), (15, 23, 42, 255))
    draw = ImageDraw.Draw(img)

    try:
        font_axis = ImageFont.truetype("C:/Windows/Fonts/segoeui.ttf", 20)
        font_legend = ImageFont.truetype("C:/Windows/Fonts/segoeuib.ttf", 20)
        font_label = ImageFont.truetype("C:/Windows/Fonts/segoeuib.ttf", 19)
    except Exception:
        font_axis = font_legend = font_label = ImageFont.load_default()

    padding_top = 100
    padding_bottom = 120
    padding_left = 120
    padding_right = 60

    chart_w = width - padding_left - padding_right
    chart_h = height - padding_top - padding_bottom

    # Card border
    draw.rounded_rectangle([10, 10, width - 10, height - 10], radius=20, outline=(51, 65, 85), width=3)

    # Series colors (100% Distinct High-Contrast Colors)
    c_tot = (168, 85, 247)   # #a855f7 Electric Purple
    c_mnj = (239, 68, 68)   # #ef4444 Crimson Red
    c_kx = (6, 182, 212)     # #06b6d4 Cyan Aqua
    c_text = (148, 163, 184) # #94a3b8

    # Top Legend
    legend_y = 45
    legends = [
        ("DOI Combined Total", c_tot),
        ("DOI MNJ (Distributor)", c_mnj),
        ("DOI KX (Principal)", c_kx)
    ]

    leg_x = padding_left
    for label, col in legends:
        draw.ellipse([leg_x, legend_y + 4, leg_x + 16, legend_y + 20], fill=col)
        draw.text((leg_x + 24, legend_y), label, fill=(241, 245, 249), font=font_legend)
        bbox = font_legend.getbbox(label)
        leg_w = bbox[2] - bbox[0]
        leg_x += leg_w + 60

    if not doi_trend:
        return img

    all_vals = []
    for d in doi_trend:
        all_vals.extend([d.get("doi_total_days", 0), d.get("doi_mnj_days", 0), d.get("doi_kx_days", 0)])

    max_val = max(max(all_vals, default=100), 100) * 1.15
    min_val = 0.0

    def get_y(val):
        return padding_top + chart_h - ((val - min_val) / (max_val - min_val)) * chart_h

    num_points = len(doi_trend)
    x_step = chart_w / max(1, num_points - 1)

    pts_tot = []
    pts_mnj = []
    pts_kx = []

    for i, d in enumerate(doi_trend):
        px = padding_left + i * x_step
        pts_tot.append((px, get_y(d.get("doi_total_days", 0)), d.get("doi_total_days", 0), d.get("period_label", "")))
        pts_mnj.append((px, get_y(d.get("doi_mnj_days", 0)), d.get("doi_mnj_days", 0)))
        pts_kx.append((px, get_y(d.get("doi_kx_days", 0)), d.get("doi_kx_days", 0)))

    # Y-axis labels (clean without box grid lines)
    for ratio in [0.0, 0.25, 0.5, 0.75, 1.0]:
        gy = padding_top + chart_h * (1.0 - ratio)
        g_val = min_val + (max_val - min_val) * ratio
        draw.text((padding_left - 80, gy - 12), f"{g_val:.0f}d", fill=c_text, font=font_axis)

    # X-axis labels (clean without vertical guide lines)
    for px, py, val, plabel in pts_tot:
        bbox = font_axis.getbbox(plabel)
        tw = bbox[2] - bbox[0]
        draw.text((px - tw / 2, height - padding_bottom + 25), plabel, fill=(203, 213, 225), font=font_axis)

    # Area under Total curve
    area_pts = [(padding_left, padding_top + chart_h)]
    area_pts.extend([(p[0], p[1]) for p in pts_tot])
    area_pts.append((pts_tot[-1][0], padding_top + chart_h))
    
    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    ov_draw = ImageDraw.Draw(overlay)
    ov_draw.polygon(area_pts, fill=(168, 85, 247, 40))
    img = Image.alpha_composite(img, overlay)
    draw = ImageDraw.Draw(img)

    def draw_thick_line(points, color, width_px=5):
        for k in range(len(points) - 1):
            p1 = (points[k][0], points[k][1])
            p2 = (points[k+1][0], points[k+1][1])
            draw.line([p1, p2], fill=color, width=width_px)

    draw_thick_line(pts_tot, c_tot, 7)
    draw_thick_line(pts_mnj, c_mnj, 5)
    draw_thick_line(pts_kx, c_kx, 5)

    # Data Point Value Labels (Clean lines without point markers)
    for i in range(num_points):
        px, py_t, val_t, _ = pts_tot[i]
        _, py_m, val_m = pts_mnj[i]
        _, py_k, val_k = pts_kx[i]

        # Total Label
        draw.text((px - 25, py_t - 32), f"{val_t:.1f}d", fill=(192, 132, 252), font=font_label)

        # Smart collision logic
        y_mnj_txt = py_m - 28
        y_kx_txt = py_k + 12

        if py_k < py_m:
            y_kx_txt = py_k - 28
            y_mnj_txt = py_m + 12

        if abs(y_mnj_txt - (py_t - 28)) < 24:
            y_mnj_txt = py_m + 12
        if abs(y_kx_txt - (py_t - 28)) < 24:
            y_kx_txt = py_k + 12

        # MNJ Label
        draw.text((px - 25, y_mnj_txt), f"{val_m:.1f}d", fill=(248, 113, 113), font=font_label)

        # KX Label
        draw.text((px - 25, y_kx_txt), f"{val_k:.1f}d", fill=(34, 211, 238), font=font_label)

    return img

def generate_doi_ppt(data_engine, filters: Dict[str, Any], template_path: str = "Template PPT.pptx") -> io.BytesIO:
    """Generates a PowerPoint presentation using Template PPT.pptx based on active dashboard filters and data."""
    if not os.path.exists(template_path):
        prs = Presentation()
    else:
        prs = Presentation(template_path)

    period = filters.get("period", "2026-07")
    unit = filters.get("unit", "value").lower()
    is_value = (unit == "value")
    avg_months = int(filters.get("avg_months", "6"))
    health_status = filters.get("health_status", "All")
    gb = filters.get("gb", "All")
    keterangan = filters.get("keterangan", "All")
    products = filters.get("products", "All")

    period_label = format_period_label(period)

    # Fetch Data from ETL engine
    summary = data_engine.get_doi_mnj_report(period=period, avg_months=avg_months)
    gb_summary = data_engine.get_gb_summary_report(period=period, avg_months=avg_months, keterangan=keterangan, unit=unit, products=products, health_status=health_status)
    doi_trend = data_engine.get_historical_doi_trend(gb=gb, keterangan=keterangan, avg_months=avg_months, unit=unit, products=products, health_status=health_status, until_period=period)

    # Filtered full report
    full_filtered = [r for r in summary if (gb == "All" or r["gb"] == gb) and (keterangan == "All" or r["keterangan_produk"] == keterangan)]

    # Calculate status scorecard counts
    total_sku = len(full_filtered)
    under_cnt = sum(1 for r in full_filtered if r["health_status_total"] == "Understock")
    norm_cnt = sum(1 for r in full_filtered if r["health_status_total"] == "Normal")
    over_cnt = sum(1 for r in full_filtered if r["health_status_total"] == "Overstock")

    # --- SLIDE 1: COVER (Populate Template Title & Subtitle Placeholders inside Red Box) ---
    slide_1 = prs.slides[0]

    title_ph = None
    sub_ph = None
    for shape in slide_1.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type == 3:  # CENTER_TITLE
                title_ph = shape
            elif shape.placeholder_format.type == 4:  # SUBTITLE
                sub_ph = shape

    if title_ph and title_ph.has_text_frame:
        tf1 = title_ph.text_frame
        tf1.word_wrap = True
        tf1.text = ""
        p0 = tf1.paragraphs[0]
        p0.text = "LAPORAN MONITORING & EVALUASI DOI PERSEDIAAN"
        p0.alignment = PP_ALIGN.CENTER
        for run in p0.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(17)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

    if sub_ph and sub_ph.has_text_frame:
        sub_ph.left = Inches(5.2)
        sub_ph.top = Inches(3.3)
        sub_ph.width = Inches(4.3)
        sub_ph.height = Inches(1.2)

        tf_sub = sub_ph.text_frame
        tf_sub.word_wrap = True
        tf_sub.text = ""

        p1 = tf_sub.paragraphs[0]
        p1.text = f"Distributor MNJ & Principal Konimex (KX)\nPeriode: {period_label}"
        p1.alignment = PP_ALIGN.CENTER
        for run in p1.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)

        p2 = tf_sub.add_paragraph()
        p2.text = f"Mode: {'Valuasi (IDR)' if is_value else 'Kuantitas (Unit)'} | GB: {gb} | Status: {health_status}"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(3)
        for run in p2.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(9.5)
            run.font.color.rgb = RGBColor(254, 243, 199)

    # --- SLIDE 2: MAIN DASHBOARD OVERVIEW (4 STATUS CARDS + 3-IN-1 TREND CHART + GB SUMMARY TABLE) ---
    slide_2 = prs.slides[1] if len(prs.slides) > 1 else prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_title(slide_2, f"🎯 Evaluasi & Ringkasan DOI Persediaan (Januari 2026 – {period_label})")

    # 4 Status Cards Grid (Fits perfectly within 10.0" width: left=0.35" to 9.65")
    cards_data = [
        {"title": "Semua Status", "count": f"{total_sku} SKU", "sub": "Total SKU Terdaftar", "color": RGBColor(0, 180, 216), "left": Inches(0.35)},
        {"title": "🔴 Understock", "count": f"{under_cnt} SKU", "sub": "< 45 Hari DOI", "color": RGBColor(239, 68, 68), "left": Inches(2.70)},
        {"title": "🟢 Normal", "count": f"{norm_cnt} SKU", "sub": "45 Hari – DOI Max", "color": RGBColor(16, 185, 129), "left": Inches(5.05)},
        {"title": "🟡 Overstock", "count": f"{over_cnt} SKU", "sub": "> DOI Max Stok", "color": RGBColor(245, 158, 11), "left": Inches(7.40)},
    ]

    for card in cards_data:
        shape = slide_2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card["left"], Inches(0.80), Inches(2.25), Inches(0.85))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(15, 23, 42)
        shape.line.color.rgb = card["color"]
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        
        p0 = tf.paragraphs[0]
        p0.text = card["title"]
        p0.font.name = "Segoe UI"
        p0.font.size = Pt(8.5)
        p0.font.bold = True
        p0.font.color.rgb = card["color"]

        p1 = tf.add_paragraph()
        p1.text = card["count"]
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(12.5)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.space_before = Pt(0)

        p2 = tf.add_paragraph()
        p2.text = card["sub"]
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(7.0)
        p2.font.color.rgb = RGBColor(148, 163, 184)

    is_gb_filtered = (gb != "All")

    if is_gb_filtered:
        # Filter GB Spesifik Terpilih (misal GB 4): Chart Trend di atas, Ringkasan GB di BAWAH chart
        hdr_box = slide_2.shapes.add_textbox(Inches(0.35), Inches(1.70), Inches(9.30), Inches(0.25))
        htf = hdr_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = f"📈 Trend Pergerakan DOI Historis 3-in-1 ({gb})"
        hp.font.name = "Segoe UI"
        hp.font.size = Pt(9.5)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(2, 132, 199)

        chart_img = create_dashboard_trend_chart(doi_trend, width=1800, height=850)
        chart_buf = io.BytesIO()
        chart_img.save(chart_buf, format="PNG")
        chart_buf.seek(0)
        slide_2.shapes.add_picture(chart_buf, Inches(0.35), Inches(1.95), Inches(9.30), Inches(2.35))

        # Tabel Ringkasan GB di BAWAH Grafik Tren
        hdr_box2 = slide_2.shapes.add_textbox(Inches(0.35), Inches(4.35), Inches(9.30), Inches(0.25))
        htf2 = hdr_box2.text_frame
        hp2 = htf2.paragraphs[0]
        hp2.text = f"🏢 Ringkasan DOI Per Group Business ({gb})"
        hp2.font.name = "Segoe UI"
        hp2.font.size = Pt(9.5)
        hp2.font.bold = True
        hp2.font.color.rgb = RGBColor(2, 132, 199)

        target_gbs = [g for g in gb_summary if g["gb"] == gb]
        if not target_gbs:
            target_gbs = gb_summary

        has_total_row = (len(target_gbs) > 1)
        num_rows_s2 = len(target_gbs) + (1 if has_total_row else 0)
        table_shape2 = slide_2.shapes.add_table(num_rows_s2 + 1, 8, Inches(0.35), Inches(4.60), Inches(9.30), Inches(0.55 if not has_total_row else 0.85))
        table2 = table_shape2.table

        gb_headers_s2 = ["GB", "SKU", "Stok Combined", "Avg Sales", "DOI Total", "DOI Max", "Selisih GB", "Status"]
        for col_idx, h_text in enumerate(gb_headers_s2):
            style_cell(table2.cell(0, col_idx), h_text, font_size=8, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, bg_color=RGBColor(15, 23, 42))

        for r_idx, g in enumerate(target_gbs, start=1):
            stok_val = g["stok_total_value" if is_value else "stok_total_qty"]
            sales_val = g["avg_sales_value" if is_value else "avg_sales_qty"]
            doi_tot = g["doi_total_days"]
            doi_max = g["doi_max_days"]
            sel_stok = g["selisih_value" if is_value else "selisih_qty"]
            sel_gb_val = (sel_stok / sales_val * 30.0) if sales_val > 0 else 0.0
            status = g["health_status_total"]

            row_vals = [
                g["gb"],
                str(g["total_sku"]),
                format_curr_or_qty(stok_val, is_value),
                format_curr_or_qty(sales_val, is_value),
                f"{doi_tot:.1f} d",
                f"{doi_max:.1f} d",
                f"+{sel_gb_val:.2f} d" if sel_gb_val > 0 else (f"{sel_gb_val:.2f} d" if sel_gb_val < 0 else "0.00 d"),
                status
            ]

            bg = RGBColor(30, 41, 59) if r_idx % 2 == 1 else RGBColor(15, 23, 42)
            for c_idx, val_text in enumerate(row_vals):
                align = PP_ALIGN.LEFT if c_idx in [0, 7] else PP_ALIGN.RIGHT
                color = RGBColor(255, 255, 255)
                if c_idx == 4: color = RGBColor(0, 242, 254)
                elif c_idx == 5: color = RGBColor(167, 243, 208)
                elif c_idx == 6: color = RGBColor(56, 189, 248) if sel_gb_val > 0 else (RGBColor(248, 113, 113) if sel_gb_val < 0 else RGBColor(148, 163, 184))
                style_cell(table2.cell(r_idx, c_idx), val_text, font_size=7.5, bold=(c_idx == 0), color=color, align=align, bg_color=bg)

        if has_total_row:
            tot_sku_s2 = sum(g["total_sku"] for g in target_gbs)
            tot_stok_s2 = sum(g["stok_total_value" if is_value else "stok_total_qty"] for g in target_gbs)
            tot_sales_s2 = sum(g["avg_sales_value" if is_value else "avg_sales_qty"] for g in target_gbs)
            tot_max_s2 = sum(g["max_value_total" if is_value else "max_qty_total"] for g in target_gbs)
            tot_sel_stok_s2 = sum(g["selisih_value" if is_value else "selisih_qty"] for g in target_gbs)

            doi_tot_s2 = (tot_stok_s2 / tot_sales_s2 * 30.0) if tot_sales_s2 > 0 else 0
            doi_max_s2 = (tot_max_s2 / tot_sales_s2 * 30.0) if tot_sales_s2 > 0 else 0
            sel_gb_s2 = (tot_sel_stok_s2 / tot_sales_s2 * 30.0) if tot_sales_s2 > 0 else 0

            cons_row_s2 = [
                "TOTAL TERPILIH",
                str(tot_sku_s2),
                format_curr_or_qty(tot_stok_s2, is_value),
                format_curr_or_qty(tot_sales_s2, is_value),
                f"{doi_tot_s2:.1f} d",
                f"{doi_max_s2:.1f} d",
                f"+{sel_gb_s2:.2f} d" if sel_gb_s2 > 0 else (f"{sel_gb_s2:.2f} d" if sel_gb_s2 < 0 else "0.00 d"),
                "Overstock" if tot_stok_s2 > tot_max_s2 else "Normal"
            ]

            c_idx_last2 = num_rows_s2
            for c_idx, val_text in enumerate(cons_row_s2):
                align = PP_ALIGN.LEFT if c_idx in [0, 7] else PP_ALIGN.RIGHT
                style_cell(table2.cell(c_idx_last2, c_idx), val_text, font_size=7.5, bold=True, color=RGBColor(0, 242, 254), align=align, bg_color=RGBColor(2, 132, 199))

    else:
        # Filter GB = All: Large full-width trend chart di Slide 2, Ringkasan GB lengkap ada di Slide berikutnya (Slide 4)
        hdr_box = slide_2.shapes.add_textbox(Inches(0.35), Inches(1.70), Inches(9.30), Inches(0.25))
        htf = hdr_box.text_frame
        hp = htf.paragraphs[0]
        hp.text = "📈 Trend Pergerakan DOI Historis Konsolidasi (3-in-1)"
        hp.font.name = "Segoe UI"
        hp.font.size = Pt(10)
        hp.font.bold = True
        hp.font.color.rgb = RGBColor(2, 132, 199)

        chart_img = create_dashboard_trend_chart(doi_trend, width=1800, height=1000)
        chart_buf = io.BytesIO()
        chart_img.save(chart_buf, format="PNG")
        chart_buf.seek(0)

        slide_2.shapes.add_picture(chart_buf, Inches(0.35), Inches(1.98), Inches(9.30), Inches(3.50))

    # --- SLIDE GENERATION & MANIFEST ---
    active_slides = [slide_1, slide_2]

    # --- SLIDE 3 (ONLY WHEN gb == "All"): RINGKASAN DOI PER GROUP BUSINESS ---
    if gb == "All":
        idx_gb_slide = len(active_slides)
        slide_gb = prs.slides[idx_gb_slide] if len(prs.slides) > idx_gb_slide else prs.slides.add_slide(prs.slide_layouts[6])
        active_slides.append(slide_gb)
        set_slide_title(slide_gb, f"🏢 Ringkasan DOI Per Group Business ({period_label})")

        num_rows = len(gb_summary)
        table_shape_gb = slide_gb.shapes.add_table(num_rows + 2, 10, Inches(0.25), Inches(0.85), Inches(9.45), Inches(3.8))
        table_gb = table_shape_gb.table

        # Explicit Column Widths to prevent awkward line wraps (e.g. TOTAL KONSOLIDASI on 1 line)
        col_widths_gb = [
            Inches(1.25),  # GB (wide enough for "TOTAL KONSOLIDASI" on 1 line)
            Inches(0.45),  # SKU
            Inches(1.15),  # Stok Combined
            Inches(1.15),  # Avg Sales/Bln
            Inches(0.70),  # DOI Total
            Inches(0.70),  # DOI Max
            Inches(0.80),  # Selisih DOI
            Inches(1.25),  # Selisih Stok
            Inches(0.70),  # DOI Net
            Inches(0.85)   # Status
        ]
        for col_idx, width in enumerate(col_widths_gb):
            table_gb.columns[col_idx].width = width

        gb_headers = ["GB", "SKU", "Stok Combined", "Avg Sales/Bln", "DOI Total", "DOI Max", "Selisih DOI", "Selisih Stok", "DOI Net", "Status"]
        for col_idx, h_text in enumerate(gb_headers):
            style_cell(table_gb.cell(0, col_idx), h_text, font_size=8, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, bg_color=RGBColor(15, 23, 42))

        tot_sku_gbs = sum(g["total_sku"] for g in gb_summary)
        tot_stok_gbs = sum(g["stok_total_value" if is_value else "stok_total_qty"] for g in gb_summary)
        tot_sales_gbs = sum(g["avg_sales_value" if is_value else "avg_sales_qty"] for g in gb_summary)
        tot_max_gbs = sum(g["max_value_total" if is_value else "max_qty_total"] for g in gb_summary)
        tot_selisih_stok_gbs = sum(g["selisih_value" if is_value else "selisih_qty"] for g in gb_summary)

        for r_idx, g in enumerate(gb_summary, start=1):
            stok_val = g["stok_total_value" if is_value else "stok_total_qty"]
            sales_val = g["avg_sales_value" if is_value else "avg_sales_qty"]
            sel_stok = g["selisih_value" if is_value else "selisih_qty"]
            doi_tot = g["doi_total_days"]
            doi_max = g["doi_max_days"]
            sel_doi = g["selisih_doi_days"]
            doi_net = g["doi_after_selisih"]
            status = g["health_status_total"]

            row_vals = [
                g["gb"],
                str(g["total_sku"]),
                format_curr_or_qty(stok_val, is_value),
                format_curr_or_qty(sales_val, is_value),
                f"{doi_tot:.1f} d",
                f"{doi_max:.1f} d",
                f"+{sel_doi:.1f} d" if sel_doi > 0 else (f"{sel_doi:.1f} d" if sel_doi < 0 else "0.0 d"),
                f"+{format_curr_or_qty(sel_stok, is_value)}" if sel_stok > 0 else (format_curr_or_qty(sel_stok, is_value) if sel_stok < 0 else "0"),
                f"{doi_net:.1f} d",
                status
            ]

            bg = RGBColor(30, 41, 59) if r_idx % 2 == 1 else RGBColor(15, 23, 42)
            for c_idx, val_text in enumerate(row_vals):
                align = PP_ALIGN.LEFT if c_idx in [0, 9] else PP_ALIGN.RIGHT
                color = RGBColor(255, 255, 255)
                if c_idx == 4: color = RGBColor(0, 242, 254)
                elif c_idx == 5: color = RGBColor(167, 243, 208)
                elif c_idx in [6, 7]: color = RGBColor(251, 191, 36) if sel_stok > 0 else (RGBColor(248, 113, 113) if sel_stok < 0 else RGBColor(148, 163, 184))
                elif c_idx == 8: color = RGBColor(167, 243, 208)
                style_cell(table_gb.cell(r_idx, c_idx), val_text, font_size=7.5, bold=(c_idx == 0), color=color, align=align, bg_color=bg)

        # Consolidated Total Row
        doi_tot_cons = (tot_stok_gbs / tot_sales_gbs * 30.0) if tot_sales_gbs > 0 else 0
        doi_max_cons = (tot_max_gbs / tot_sales_gbs * 30.0) if tot_sales_gbs > 0 else 0
        sel_doi_cons = (tot_selisih_stok_gbs / tot_sales_gbs * 30.0) if tot_sales_gbs > 0 else 0
        doi_net_cons = doi_tot_cons - sel_doi_cons

        cons_row = [
            "TOTAL KONSOLIDASI",
            str(tot_sku_gbs),
            format_curr_or_qty(tot_stok_gbs, is_value),
            format_curr_or_qty(tot_sales_gbs, is_value),
            f"{doi_tot_cons:.1f} d",
            f"{doi_max_cons:.1f} d",
            f"+{sel_doi_cons:.1f} d" if sel_doi_cons > 0 else (f"{sel_doi_cons:.1f} d" if sel_doi_cons < 0 else "0.0 d"),
            f"+{format_curr_or_qty(tot_selisih_stok_gbs, is_value)}" if tot_selisih_stok_gbs > 0 else format_curr_or_qty(tot_selisih_stok_gbs, is_value),
            f"{doi_net_cons:.1f} d",
            "Overstock" if tot_stok_gbs > tot_max_gbs else "Normal"
        ]

        c_idx_last = num_rows + 1
        for c_idx, val_text in enumerate(cons_row):
            align = PP_ALIGN.LEFT if c_idx in [0, 9] else PP_ALIGN.RIGHT
            style_cell(table_gb.cell(c_idx_last, c_idx), val_text, font_size=7.5, bold=True, color=RGBColor(0, 242, 254), align=align, bg_color=RGBColor(2, 132, 199))

    # Calculate selisih_doi_gb for each item in summary
    total_avg_sales_all = sum(r.get("avg_sales_value", 0.0) for r in summary)
    gb_sales_dict = {g["gb"]: g.get("avg_sales_value", 0.0) for g in gb_summary}

    for item in full_filtered:
        sel_val = item.get("selisih_value", 0.0)
        denom = total_avg_sales_all if gb == "All" else gb_sales_dict.get(item["gb"], 0.0)
        if denom > 0:
            item["selisih_doi_gb"] = (sel_val / denom) * 30.0
        else:
            item["selisih_doi_gb"] = 0.0

    items_per_page = 11

    # --- 1. OVERSTOCK DETAIL SLIDES (TOP 20 PARETO: SELISIH QTY, SELISIH DOI, SELISIH GB) ---
    overstock_items = [r for r in full_filtered if r.get("health_status_total") == "Overstock"]
    # Sort Pareto: highest selisih_doi_gb / selisih_value descending
    overstock_items.sort(key=lambda x: (x.get("selisih_doi_gb", 0.0), x.get("selisih_value", 0.0)), reverse=True)
    # Take Top 20 items ONLY!
    overstock_items = overstock_items[:20]

    items_per_page_over = 10
    chunked_over = [overstock_items[i:i + items_per_page_over] for i in range(0, len(overstock_items), items_per_page_over)]
    if not chunked_over:
        chunked_over = [[]]

    for chunk_idx, chunk in enumerate(chunked_over):
        idx_slide = len(active_slides)
        slide_over = prs.slides[idx_slide] if len(prs.slides) > idx_slide else prs.slides.add_slide(prs.slide_layouts[6])
        active_slides.append(slide_over)

        page_suffix = f" (Hal {chunk_idx + 1}/{len(chunked_over)})" if len(chunked_over) > 1 else ""
        gb_title_suffix = f" ({gb})" if gb != "All" else ""
        set_slide_title(slide_over, f"🟡 Top 20 SKU Overstock – Pareto Selisih{gb_title_suffix}{page_suffix}")

        num_rows = len(chunk)
        table_shape = slide_over.shapes.add_table(num_rows + 1, 11, Inches(0.20), Inches(0.85), Inches(9.55), Inches(3.8))
        table = table_shape.table

        col_widths = [
            Inches(0.75),  # Kode
            Inches(2.10),  # Nama Produk
            Inches(0.45),  # GB
            Inches(0.65),  # Ket
            Inches(0.85),  # Stok (Qty)
            Inches(0.85),  # Sales (Qty)
            Inches(0.65),  # DOI Total
            Inches(0.65),  # DOI Max
            Inches(0.95),  # Selisih Qty
            Inches(0.80),  # Selisih DOI
            Inches(0.80)   # Selisih GB
        ]
        for col_idx, width in enumerate(col_widths):
            table.columns[col_idx].width = width

        headers = ["Kode", "Nama Produk", "GB", "Ket", "Stok (Qty)", "Sales (Qty)", "DOI Total", "DOI Max", "Selisih Qty", "Selisih DOI", "Selisih GB"]
        for col_idx, h_text in enumerate(headers):
            style_cell(table.cell(0, col_idx), h_text, font_size=8, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, bg_color=RGBColor(15, 23, 42))

        for r_idx, item in enumerate(chunk, start=1):
            stok_qty = item.get("stok_total_qty", 0)
            sales_qty = item.get("avg_sales_qty", 0)
            doi_tot = item.get("doi_total_days", 0.0)
            doi_max = item.get("doi_max_days", 90.0)
            sel_qty = item.get("selisih_qty", 0)
            sel_doi = item.get("selisih_doi_days", 0.0)
            sel_gb_val = item.get("selisih_doi_gb", 0.0)

            row_vals = [
                item.get("product_code", "-"),
                item.get("product_name", "-")[:32],
                item.get("gb", "-"),
                item.get("keterangan_produk", "-"),
                f"{int(stok_qty):,}",
                f"{int(sales_qty):,}",
                f"{doi_tot:.1f} d",
                f"{doi_max:.1f} d",
                f"+{int(sel_qty):,}" if sel_qty > 0 else f"{int(sel_qty):,}",
                f"+{sel_doi:.1f} d" if sel_doi > 0 else f"{sel_doi:.1f} d",
                f"+{sel_gb_val:.2f} d" if sel_gb_val > 0 else f"{sel_gb_val:.2f} d"
            ]

            bg = RGBColor(30, 41, 59) if r_idx % 2 == 1 else RGBColor(15, 23, 42)
            for c_idx, val_text in enumerate(row_vals):
                align = PP_ALIGN.LEFT if c_idx in [0, 1, 2, 3] else PP_ALIGN.RIGHT
                color = RGBColor(255, 255, 255)
                if c_idx == 6: color = RGBColor(0, 242, 254)
                elif c_idx == 7: color = RGBColor(167, 243, 208)
                elif c_idx in [8, 9, 10]: color = RGBColor(251, 191, 36)
                style_cell(table.cell(r_idx, c_idx), val_text, font_size=7.5, bold=(c_idx in [0, 8, 10]), color=color, align=align, bg_color=bg)

    # --- 2. UNDERSTOCK DETAIL SLIDES (ALL UNDERSTOCK ITEMS, PAGINATED 10 ITEMS PER SLIDE) ---
    understock_items = [r for r in full_filtered if r.get("health_status_total") == "Understock"]
    # Sort highest deficit impact first
    understock_items.sort(key=lambda x: abs(x.get("selisih_value", 0.0)), reverse=True)

    items_per_page_under = 10
    chunked_under = [understock_items[i:i + items_per_page_under] for i in range(0, len(understock_items), items_per_page_under)]
    if not chunked_under:
        chunked_under = [[]]

    for chunk_idx, chunk in enumerate(chunked_under):
        idx_slide = len(active_slides)
        slide_under = prs.slides[idx_slide] if len(prs.slides) > idx_slide else prs.slides.add_slide(prs.slide_layouts[6])
        active_slides.append(slide_under)

        page_suffix = f" (Hal {chunk_idx + 1}/{len(chunked_under)})" if len(chunked_under) > 1 else ""
        gb_title_suffix = f" ({gb})" if gb != "All" else ""
        set_slide_title(slide_under, f"🔴 Detail SKU Understock – Defisit Stok Qty{gb_title_suffix}{page_suffix}")

        num_rows = len(chunk)
        table_shape = slide_under.shapes.add_table(num_rows + 1, 10, Inches(0.20), Inches(0.85), Inches(9.55), Inches(3.80))
        table = table_shape.table

        col_widths = [
            Inches(0.75),  # Kode
            Inches(2.10),  # Nama Produk
            Inches(0.45),  # GB
            Inches(0.65),  # Ket
            Inches(0.95),  # Stok (Qty)
            Inches(0.95),  # Sales (Qty)
            Inches(0.75),  # DOI Total
            Inches(0.75),  # DOI Max
            Inches(1.10),  # Defisit Qty
            Inches(1.10)   # Status
        ]
        for col_idx, width in enumerate(col_widths):
            table.columns[col_idx].width = width

        # Explicit row heights so that even when len(chunk) < 10, row height and table box look clean & matching
        for row in table.rows:
            row.height = Inches(0.32)

        headers = ["Kode", "Nama Produk", "GB", "Ket", "Stok (Qty)", "Sales (Qty)", "DOI Total", "DOI Max", "Defisit Qty", "Status"]
        for col_idx, h_text in enumerate(headers):
            style_cell(table.cell(0, col_idx), h_text, font_size=8, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, bg_color=RGBColor(15, 23, 42))

        for r_idx, item in enumerate(chunk, start=1):
            stok_qty = item.get("stok_total_qty", 0)
            sales_qty = item.get("avg_sales_qty", 0)
            doi_tot = item.get("doi_total_days", 0.0)
            doi_max = item.get("doi_max_days", 90.0)
            sel_qty = item.get("selisih_qty", 0)

            row_vals = [
                item.get("product_code", "-"),
                item.get("product_name", "-")[:32],
                item.get("gb", "-"),
                item.get("keterangan_produk", "-"),
                f"{int(stok_qty):,}",
                f"{int(sales_qty):,}",
                f"{doi_tot:.1f} d",
                f"{doi_max:.1f} d",
                f"{int(sel_qty):,}",
                "Understock"
            ]

            bg = RGBColor(30, 41, 59) if r_idx % 2 == 1 else RGBColor(15, 23, 42)
            for c_idx, val_text in enumerate(row_vals):
                align = PP_ALIGN.LEFT if c_idx in [0, 1, 2, 3, 9] else PP_ALIGN.RIGHT
                color = RGBColor(255, 255, 255)
                if c_idx == 6: color = RGBColor(0, 242, 254)
                elif c_idx == 7: color = RGBColor(167, 243, 208)
                elif c_idx in [8, 9]: color = RGBColor(248, 113, 113)
                style_cell(table.cell(r_idx, c_idx), val_text, font_size=7.5, bold=(c_idx in [0, 8, 9]), color=color, align=align, bg_color=bg)

    # Clean up extra slides in prs template so presentation contains ONLY active_slides!
    keep_count = len(active_slides)
    for i in range(len(prs.slides) - 1, keep_count - 1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # Save to memory buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
